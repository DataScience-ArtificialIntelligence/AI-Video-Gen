# from fastapi import FastAPI, HTTPException
# from fastapi.responses import FileResponse
# from fastapi.middleware.cors import CORSMiddleware  # <-- add this
# from pydantic import BaseModel
# from pptx import Presentation
# from pptx.util import Inches, Pt
# from pptx.dml.color import RGBColor
# import uuid
# import os

# app = FastAPI()

# # CORS settings
# origins = [
#     "http://localhost:5173",  # React Vite default
#     "http://127.0.0.1:5173"
# ]

# app.add_middleware(
#     CORSMiddleware,
#     allow_origins=origins,
#     allow_credentials=True,
#     allow_methods=["*"],   # allow all HTTP methods
#     allow_headers=["*"],   # allow all headers
# )

# class PromptRequest(BaseModel):
#     prompt: str

# def create_ppt(prompt: str, file_path: str):
#     prs = Presentation()
    
#     # Title Slide
#     slide_layout = prs.slide_layouts[0]
#     slide = prs.slides.add_slide(slide_layout)
#     slide.shapes.title.text = "Generated PPT"
#     slide.placeholders[1].text = prompt

#     # Add 3 content slides
#     for i in range(1, 4):
#         slide_layout = prs.slide_layouts[1]
#         slide = prs.slides.add_slide(slide_layout)
#         slide.shapes.title.text = f"Slide {i}"
#         slide.placeholders[1].text = f"Content related to '{prompt}' on slide {i}"

#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 for paragraph in shape.text_frame.paragraphs:
#                     for run in paragraph.runs:
#                         run.font.color.rgb = RGBColor(i*50, 100, 200-i*50)
#                         run.font.size = Pt(24)

#     prs.save(file_path)

# @app.post("/generate-ppt")
# async def generate_ppt(request: PromptRequest):
#     try:
#         file_name = f"{uuid.uuid4()}.pptx"
#         file_path = f"./{file_name}"
#         create_ppt(request.prompt, file_path)
#         return FileResponse(
#             file_path,
#             filename="GeneratedPPT.pptx",
#             media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
#         )
#     except Exception as e:
#         raise HTTPException(status_code=500, detail=str(e))
from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import uuid
import random

app = FastAPI()

# CORS settings
origins = [
    "http://localhost:5173",  # Vite React default
    "http://127.0.0.1:5173"
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=origins,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class PromptRequest(BaseModel):
    prompt: str

def create_ppt(prompt: str, file_path: str, num_slides=5):
    prs = Presentation()

    # Predefined colors
    bg_colors = [
        RGBColor(255, 179, 186),
        RGBColor(255, 223, 186),
        RGBColor(255, 255, 186),
        RGBColor(186, 255, 201),
        RGBColor(186, 225, 255),
    ]
    text_colors = [
        RGBColor(50, 50, 50),
        RGBColor(0, 102, 204),
        RGBColor(102, 0, 51),
        RGBColor(0, 153, 0),
        RGBColor(128, 0, 128),
    ]

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = prompt
    slide.placeholders[1].text = f"A presentation about {prompt}"
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = random.choice(bg_colors)

    # Content Slides
    for i in range(1, num_slides):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = f"Slide {i}: {prompt.split()[0].capitalize()} Topic"
        slide.placeholders[1].text = f"This slide covers key points of '{prompt}' (point {i})."

        # Random background
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = random.choice(bg_colors)

        # Colorful text
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = random.choice(text_colors)
                        run.font.size = Pt(random.choice([20, 22, 24, 26]))

    prs.save(file_path)

@app.post("/generate-ppt")
async def generate_ppt(request: PromptRequest):
    try:
        file_name = f"{uuid.uuid4()}.pptx"
        file_path = f"./{file_name}"
        create_ppt(request.prompt, file_path)
        return FileResponse(
            file_path,
            filename="GeneratedPPT.pptx",
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
